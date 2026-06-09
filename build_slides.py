"""Generate PRESENTATION.pptx from the project figures (non-technical deck)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette ---------------------------------------------------------------
NAVY   = RGBColor(0x0B, 0x2A, 0x4A)   # deep blue
BLUE   = RGBColor(0x1F, 0x6F, 0xB2)   # accent blue
LIGHT  = RGBColor(0xED, 0xF3, 0xF9)   # pale panel
GREY   = RGBColor(0x55, 0x5B, 0x66)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GOLD   = RGBColor(0xC8, 0x96, 0x2C)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold)."""
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for (t, sz, col, bold) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            r.font.name = "Calibri"
    return tb


def header(s, kicker, title):
    rect(s, 0, 0, SW, Inches(1.35), NAVY)
    rect(s, 0, Inches(1.35), SW, Inches(0.06), GOLD)
    txt(s, Inches(0.6), Inches(0.18), Inches(12), Inches(0.4),
        [[(kicker.upper(), 13, RGBColor(0x9F,0xC4,0xE4), True)]])
    txt(s, Inches(0.6), Inches(0.5), Inches(12.1), Inches(0.8),
        [[(title, 30, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)


def footer(s, n):
    txt(s, Inches(0.6), Inches(7.05), Inches(9), Inches(0.35),
        [[("Corpus des banques centrales  ·  document de présentation  ·  juin 2026", 9, GREY, False)]])
    txt(s, Inches(12.2), Inches(7.05), Inches(0.8), Inches(0.35),
        [[(str(n), 9, GREY, True)]], align=PP_ALIGN.RIGHT)


def chip(s, x, y, w, h, big, label, col=BLUE):
    rect(s, x, y, w, h, LIGHT)
    rect(s, x, y, Inches(0.08), h, col)
    txt(s, x+Inches(0.22), y+Inches(0.12), w-Inches(0.35), Inches(0.7),
        [[(big, 30, NAVY, True)]])
    txt(s, x+Inches(0.22), y+h-Inches(0.55), w-Inches(0.35), Inches(0.5),
        [[(label, 12.5, GREY, False)]])


# ===========================================================================
# 1. TITLE
# ===========================================================================
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.55), SW, Inches(0.08), GOLD)
txt(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.6),
    [[("La bibliothèque mondiale", 46, WHITE, True)],
     [("des banques centrales", 46, RGBColor(0x9F,0xC4,0xE4), True)]],
    space_after=2)
txt(s, Inches(0.95), Inches(4.8), Inches(11.5), Inches(1.2),
    [[("Rassembler tous les documents officiels publiés par les banques "
       "centrales du monde — uniquement des sources de première main.", 18,
       RGBColor(0xD7,0xE3,0xEF), False)]])
txt(s, Inches(0.95), Inches(6.6), Inches(11.5), Inches(0.5),
    [[("38 413 documents  ·  61 banques centrales  ·  1969 – 2026", 15, GOLD, True)]])

# ===========================================================================
# 2. L'IDÉE
# ===========================================================================
s = slide(); header(s, "Le projet", "L'idée, en clair")
txt(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(1.3),
    [[("Les banques centrales publient en permanence — discours, rapports, "
       "décisions de taux, minutes de réunion. Cette information est ", 17, GREY, False),
      ("publique mais éparpillée", 17, NAVY, True),
      (" : chaque banque a son site, son format, sa langue.", 17, GREY, False)],
     [("Nous rassemblons tout cela ", 17, GREY, False),
      ("automatiquement et au même endroit", 17, NAVY, True),
      (", pour pouvoir l'interroger en langage naturel et répondre avec les "
       "documents officiels eux-mêmes, source à l'appui.", 17, GREY, False)]],
    line_spacing=1.08, space_after=10)

cards = [
    ("Uniquement officiel",
     "Chaque fichier vient du site de la banque émettrice (ou de la BIS / de "
     "l'éditeur de l'article). Aucune source non officielle."),
    ("Aucune analyse",
     "Pas de presse, pas de notes d'analystes, rien généré par une IA. "
     "Seulement la parole brute des institutions."),
    ("Viser l'exhaustivité",
     "Couvrir toutes les grandes banques et tous les types de documents. "
     "Ce qui manque est identifié et planifié."),
]
x = Inches(0.6); w = Inches(3.9); gap = Inches(0.2)
for i, (h, body) in enumerate(cards):
    cx = x + i*(w+gap)
    rect(s, cx, Inches(4.1), w, Inches(2.5), LIGHT)
    rect(s, cx, Inches(4.1), w, Inches(0.12), GOLD)
    txt(s, cx+Inches(0.25), Inches(4.35), w-Inches(0.5), Inches(0.6),
        [[(f"{i+1}.  {h}", 17, NAVY, True)]])
    txt(s, cx+Inches(0.25), Inches(4.95), w-Inches(0.5), Inches(1.5),
        [[(body, 13.5, GREY, False)]], line_spacing=1.05)
footer(s, 2)

# ===========================================================================
# 3. CE QU'ON A DÉJÀ
# ===========================================================================
s = slide(); header(s, "État des lieux", "Ce que nous avons déjà collecté")
stats = [("38 413", "documents officiels"),
         ("61", "banques centrales"),
         ("1969-2026", "~57 ans couverts"),
         ("21 Go", "de données traçables")]
x = Inches(0.6); w = Inches(2.95); gap = Inches(0.13)
for i, (b, l) in enumerate(stats):
    chip(s, x+i*(w+gap), Inches(1.7), w, Inches(1.4), b, l)

txt(s, Inches(0.6), Inches(3.4), Inches(12), Inches(0.5),
    [[("Répartition par type de document", 18, NAVY, True)]])
rows = [("Discours & interviews des dirigeants", "18 389"),
        ("Articles de recherche (working + discussion)", "16 380"),
        ("Décisions de taux & déclarations de politique", "1 064"),
        ("Minutes / comptes-rendus de réunion", "808"),
        ("Rapports : MP, FSR, annuel, convergence…", "1 032"),
        ("Bulletins éco. (ECB + Beige Book Fed)", "485"),
        ("Projections / prévisions chiffrées", "74")]
y = Inches(3.95); rh = Inches(0.42)
for i, (lab, val) in enumerate(rows):
    bg = LIGHT if i % 2 == 0 else WHITE
    rect(s, Inches(0.6), y, Inches(12.1), rh, bg)
    txt(s, Inches(0.85), y+Inches(0.05), Inches(10), Inches(0.4),
        [[(lab, 14, NAVY, False)]])
    txt(s, Inches(10.8), y+Inches(0.07), Inches(1.7), Inches(0.45),
        [[(val, 16, BLUE, True)]], align=PP_ALIGN.RIGHT)
    y += rh
footer(s, 3)

# ===========================================================================
# 4. LES PAYS — TOP
# ===========================================================================
s = slide(); header(s, "Couverture", "61 banques centrales, du monde entier")
top = [("Banque centrale européenne", "8 599"), ("Réserve fédérale — États-Unis", "7 000"),
       ("Bank of England", "2 486"), ("Reserve Bank of Australia", "2 052"),
       ("Bank of Canada", "2 044"), ("Banca d'Italia", "1 675"),
       ("Deutsche Bundesbank", "1 646"), ("Banque de France", "1 467"),
       ("Banco de España", "1 386"), ("Bank of Japan", "1 257")]
txt(s, Inches(0.6), Inches(1.6), Inches(6), Inches(0.4),
    [[("Les mieux couvertes", 16, NAVY, True)]])
y = Inches(2.1); rh = Inches(0.46)
for i, (lab, val) in enumerate(top):
    bg = LIGHT if i % 2 == 0 else WHITE
    rect(s, Inches(0.6), y, Inches(6.2), rh, bg)
    txt(s, Inches(0.78), y+Inches(0.05), Inches(5), Inches(0.4), [[(lab, 12.5, NAVY, False)]])
    txt(s, Inches(5.7), y+Inches(0.05), Inches(1.0), Inches(0.4),
        [[(val, 13, BLUE, True)]], align=PP_ALIGN.RIGHT)
    y += rh

rect(s, Inches(7.1), Inches(2.1), Inches(5.6), Inches(4.6), LIGHT)
txt(s, Inches(7.35), Inches(2.25), Inches(5.2), Inches(0.4),
    [[("Également présentes", 16, NAVY, True)]])
txt(s, Inches(7.35), Inches(2.75), Inches(5.2), Inches(2.6),
    [[("Inde, Suède, Suisse, Malaisie, Philippines, Pays-Bas, Afrique du Sud, "
       "Irlande, Singapour, Norvège, Hong Kong, Thaïlande, Finlande, "
       "Nouvelle-Zélande, Grèce, Chine, Serbie, Chili, Turquie, Corée du Sud, "
       "Israël, Mexique, Macédoine du Nord, Islande, Portugal, Roumanie, "
       "Danemark, Tchéquie, Bulgarie… et d'autres.", 12.5, GREY, False)]],
    line_spacing=1.12)
txt(s, Inches(7.35), Inches(5.6), Inches(5.2), Inches(1.0),
    [[("À noter : ", 12.5, GOLD, True),
      ("les working papers concentrent le volume sur ~13 grandes banques ; "
       "les plus petites n'ont (pour l'instant) que des discours.", 12.5, GREY, False)]],
    line_spacing=1.1)
footer(s, 4)

# ===========================================================================
# 5. LES 5 PAYS ABSENTS
# ===========================================================================
s = slide(); header(s, "Couverture", "Les 2 pays encore absents")
txt(s, Inches(0.6), Inches(1.65), Inches(12), Inches(0.7),
    [[("61 banques sur les 63 visées sont présentes. Les 3 « absents » précédents "
       "(Lettonie, Lituanie, Luxembourg) ont été récupérés — leur absence était un "
       "bug d'attribution, corrigé. Restent 2 absences réelles :", 16, GREY, False)]])
miss = [("Pérou — Central Reserve Bank", "Publie essentiellement en espagnol — pas de version EN indexée par la BIS"),
        ("Vietnam — State Bank", "Publie essentiellement en vietnamien — idem")]
y = Inches(2.9); rh = Inches(0.9)
for i, (lab, why) in enumerate(miss):
    rect(s, Inches(0.6), y, Inches(12.1), rh-Inches(0.12), LIGHT)
    rect(s, Inches(0.6), y, Inches(0.1), rh-Inches(0.12), GOLD)
    txt(s, Inches(0.9), y+Inches(0.06), Inches(4.6), Inches(0.6),
        [[(lab, 15, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.6), y+Inches(0.06), Inches(6.9), Inches(0.6),
        [[(why, 13.5, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += rh
footer(s, 5)

# ===========================================================================
# 6. FEUILLE DE ROUTE
# ===========================================================================
s = slide(); header(s, "La suite", "Ce qu'il reste à collecter")
txt(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.5),
    [[("✅ Discours et articles de recherche sont collectés. Restent trois "
       "chantiers, du plus prioritaire au plus complexe :", 15, GREY, False)]])
items = [
    ("1", "Décisions — grandes banques faites", GOLD,
     "Décisions de taux, déclarations & minutes : 1 872 collectées "
     "(BCE, Fed, Australie, Bank of England).",
     "Largement fait pour les grandes banques ; extension site par site au reste."),
    ("2", "Conférences de presse", BLUE,
     "Les transcriptions : le « pourquoi » derrière chaque décision, "
     "questions des journalistes incluses.",
     "Le plus exigeant techniquement. À partir de la Fed et la BCE."),
    ("3", "Stabilité financière & rapports annuels", NAVY,
     "Rapports de stabilité financière et rapports annuels des banques.",
     "Patterns simples à ajouter par banque."),
]
y = Inches(2.35); rh = Inches(1.35)
for num, title, col, body, status in items:
    rect(s, Inches(0.6), y, Inches(12.1), rh-Inches(0.13), LIGHT)
    rect(s, Inches(0.6), y, Inches(1.0), rh-Inches(0.13), col)
    txt(s, Inches(0.6), y, Inches(1.0), rh-Inches(0.13),
        [[(num, 34, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.8), y+Inches(0.12), Inches(10.7), Inches(0.4),
        [[(title, 17, NAVY, True)]])
    txt(s, Inches(1.8), y+Inches(0.52), Inches(10.7), Inches(0.5),
        [[(body, 13, GREY, False)]], line_spacing=1.0)
    txt(s, Inches(1.8), y+rh-Inches(0.42), Inches(10.7), Inches(0.35),
        [[("Statut : ", 12, col, True), (status, 12, GREY, False)]])
    y += rh
footer(s, 6)

# ===========================================================================
# 7. RÉSUMÉ
# ===========================================================================
s = slide(); header(s, "Synthèse", "Aujourd'hui → objectif")
cols = [("", "Aujourd'hui", "Objectif"),
        ("Documents", "38 413", "45 000+"),
        ("Pays couverts", "61 / 63", "63 / 63"),
        ("Types de documents", "discours + recherche + décisions + minutes + rapports", "+ conférences de presse"),
        ("Nature des sources", "100 % officiel", "inchangé — c'est le principe")]
y = Inches(1.8); rh = Inches(0.72)
cw = [Inches(3.6), Inches(4.25), Inches(4.25)]
cx0 = Inches(0.6)
for i, row in enumerate(cols):
    header_row = (i == 0)
    cx = cx0
    for j, cell in enumerate(row):
        if header_row:
            rect(s, cx, y, cw[j], rh, NAVY if j == 0 else (BLUE if j == 1 else GOLD))
            color = WHITE
        else:
            rect(s, cx, y, cw[j], rh, LIGHT if i % 2 else WHITE)
            color = NAVY if j == 0 else GREY
        bold = header_row or j == 0
        txt(s, cx+Inches(0.2), y, cw[j]-Inches(0.3), rh,
            [[(cell, 14.5 if not header_row else 15, color, bold)]],
            anchor=MSO_ANCHOR.MIDDLE)
        cx += cw[j]
    y += rh

rect(s, Inches(0.6), y+Inches(0.2), Inches(12.1), Inches(1.15), NAVY)
txt(s, Inches(0.9), y+Inches(0.35), Inches(11.5), Inches(0.9),
    [[("Ce qui ne changera jamais : ", 15, GOLD, True),
      ("uniquement des sources officielles de première main, traçables, "
       "sans aucun commentaire ni contenu généré. Une base fiable, à "
       "vocation exhaustive, sur la parole des banques centrales du monde.",
       15, WHITE, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.08)
footer(s, 7)

prs.save("PRESENTATION.pptx")
print("OK — PRESENTATION.pptx,", len(prs.slides._sldIdLst), "slides")
