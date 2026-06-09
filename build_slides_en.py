"""Generate PRESENTATION_EN.pptx from the project figures (non-technical deck, English)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette ---------------------------------------------------------------
NAVY   = RGBColor(0x0B, 0x2A, 0x4A)   # deep blue
BLUE   = RGBColor(0x1F, 0x6F, 0xB2)   # accent blue
LIGHT  = RGBColor(0xED, 0xF3, 0xF9)   # pale panel
GREY   = RGBColor(0x55, 0x5B, 0x66)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GOLD   = RGBColor(0xC8, 0x96, 0x2C)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold)."""
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for (t, sz, col, bold) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            r.font.name = "Calibri"
    return tb


def header(s, kicker, title):
    rect(s, 0, 0, SW, Inches(1.35), NAVY)
    rect(s, 0, Inches(1.35), SW, Inches(0.06), GOLD)
    txt(s, Inches(0.6), Inches(0.18), Inches(12), Inches(0.4),
        [[(kicker.upper(), 13, RGBColor(0x9F,0xC4,0xE4), True)]])
    txt(s, Inches(0.6), Inches(0.5), Inches(12.1), Inches(0.8),
        [[(title, 30, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)


def footer(s, n):
    txt(s, Inches(0.6), Inches(7.05), Inches(9), Inches(0.35),
        [[("Central bank corpus  ·  briefing deck  ·  June 2026", 9, GREY, False)]])
    txt(s, Inches(12.2), Inches(7.05), Inches(0.8), Inches(0.35),
        [[(str(n), 9, GREY, True)]], align=PP_ALIGN.RIGHT)


def chip(s, x, y, w, h, big, label, col=BLUE):
    rect(s, x, y, w, h, LIGHT)
    rect(s, x, y, Inches(0.08), h, col)
    txt(s, x+Inches(0.22), y+Inches(0.12), w-Inches(0.35), Inches(0.7),
        [[(big, 30, NAVY, True)]])
    txt(s, x+Inches(0.22), y+h-Inches(0.55), w-Inches(0.35), Inches(0.5),
        [[(label, 12.5, GREY, False)]])


# ===========================================================================
# 1. TITLE
# ===========================================================================
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.55), SW, Inches(0.08), GOLD)
txt(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.6),
    [[("The world's central bank", 46, WHITE, True)],
     [("document library", 46, RGBColor(0x9F,0xC4,0xE4), True)]],
    space_after=2)
txt(s, Inches(0.95), Inches(4.8), Inches(11.5), Inches(1.2),
    [[("Bringing together every official document published by the world's "
       "central banks — first-hand sources only.", 18,
       RGBColor(0xD7,0xE3,0xEF), False)]])
txt(s, Inches(0.95), Inches(6.6), Inches(11.5), Inches(0.5),
    [[("38,413 documents  ·  61 central banks  ·  1969 – 2026", 15, GOLD, True)]])

# ===========================================================================
# 2. THE IDEA
# ===========================================================================
s = slide(); header(s, "The project", "The idea, in plain terms")
txt(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(1.3),
    [[("Central banks publish constantly — speeches, reports, rate decisions, "
       "meeting minutes. That information is ", 17, GREY, False),
      ("public but scattered", 17, NAVY, True),
      (": every bank has its own website, format and language.", 17, GREY, False)],
     [("We gather all of it ", 17, GREY, False),
      ("automatically, in one place", 17, NAVY, True),
      (", so it can be queried in plain language and answered with the "
       "official documents themselves, source included.", 17, GREY, False)]],
    line_spacing=1.08, space_after=10)

cards = [
    ("Official only",
     "Every file comes from the issuing bank's own site (or the BIS / the paper's "
     "publisher). No non-official source."),
    ("No analysis",
     "No press, no analyst notes, nothing AI-generated. Only the raw words "
     "of the institutions themselves."),
    ("Built to be exhaustive",
     "Cover every major bank and every document type. Whatever is missing "
     "is identified and planned."),
]
x = Inches(0.6); w = Inches(3.9); gap = Inches(0.2)
for i, (h, body) in enumerate(cards):
    cx = x + i*(w+gap)
    rect(s, cx, Inches(4.1), w, Inches(2.5), LIGHT)
    rect(s, cx, Inches(4.1), w, Inches(0.12), GOLD)
    txt(s, cx+Inches(0.25), Inches(4.35), w-Inches(0.5), Inches(0.6),
        [[(f"{i+1}.  {h}", 17, NAVY, True)]])
    txt(s, cx+Inches(0.25), Inches(4.95), w-Inches(0.5), Inches(1.5),
        [[(body, 13.5, GREY, False)]], line_spacing=1.05)
footer(s, 2)

# ===========================================================================
# 3. WHAT WE HAVE
# ===========================================================================
s = slide(); header(s, "Status", "What we have already collected")
stats = [("38,413", "official documents"),
         ("61", "central banks"),
         ("1969-2026", "~57 years covered"),
         ("21 GB", "of traceable data")]
x = Inches(0.6); w = Inches(2.95); gap = Inches(0.13)
for i, (b, l) in enumerate(stats):
    chip(s, x+i*(w+gap), Inches(1.7), w, Inches(1.4), b, l)

txt(s, Inches(0.6), Inches(3.4), Inches(12), Inches(0.5),
    [[("Breakdown by document type", 18, NAVY, True)]])
rows = [("Speeches & interviews of senior officials", "18,389"),
        ("Research papers (working + discussion)", "16,380"),
        ("Rate decisions & policy statements", "1,064"),
        ("Meeting minutes / accounts", "808"),
        ("Reports: MP, FSR, annual, convergence…", "1,032"),
        ("Economic bulletins (ECB + Fed Beige Book)", "485"),
        ("Projections / staff forecasts", "74")]
y = Inches(3.95); rh = Inches(0.42)
for i, (lab, val) in enumerate(rows):
    bg = LIGHT if i % 2 == 0 else WHITE
    rect(s, Inches(0.6), y, Inches(12.1), rh, bg)
    txt(s, Inches(0.85), y+Inches(0.05), Inches(10), Inches(0.4),
        [[(lab, 14, NAVY, False)]])
    txt(s, Inches(10.8), y+Inches(0.07), Inches(1.7), Inches(0.45),
        [[(val, 16, BLUE, True)]], align=PP_ALIGN.RIGHT)
    y += rh
footer(s, 3)

# ===========================================================================
# 4. THE COUNTRIES — TOP
# ===========================================================================
s = slide(); header(s, "Coverage", "61 central banks, worldwide")
top = [("European Central Bank", "8,599"), ("Federal Reserve — United States", "7,000"),
       ("Bank of England", "2,486"), ("Reserve Bank of Australia", "2,052"),
       ("Bank of Canada", "2,044"), ("Bank of Italy", "1,675"),
       ("Deutsche Bundesbank", "1,646"), ("Banque de France", "1,467"),
       ("Banco de España", "1,386"), ("Bank of Japan", "1,257")]
txt(s, Inches(0.6), Inches(1.6), Inches(6), Inches(0.4),
    [[("Best covered", 16, NAVY, True)]])
y = Inches(2.1); rh = Inches(0.46)
for i, (lab, val) in enumerate(top):
    bg = LIGHT if i % 2 == 0 else WHITE
    rect(s, Inches(0.6), y, Inches(6.2), rh, bg)
    txt(s, Inches(0.78), y+Inches(0.05), Inches(5), Inches(0.4), [[(lab, 12.5, NAVY, False)]])
    txt(s, Inches(5.7), y+Inches(0.05), Inches(1.0), Inches(0.4),
        [[(val, 13, BLUE, True)]], align=PP_ALIGN.RIGHT)
    y += rh

rect(s, Inches(7.1), Inches(2.1), Inches(5.6), Inches(4.6), LIGHT)
txt(s, Inches(7.35), Inches(2.25), Inches(5.2), Inches(0.4),
    [[("Also present", 16, NAVY, True)]])
txt(s, Inches(7.35), Inches(2.75), Inches(5.2), Inches(2.6),
    [[("India, Sweden, Switzerland, Malaysia, Philippines, Netherlands, "
       "South Africa, Ireland, Singapore, Norway, Hong Kong, Thailand, Finland, "
       "New Zealand, Greece, China, Serbia, Chile, Turkey, South Korea, Israel, "
       "Mexico, North Macedonia, Iceland, Portugal, Romania, Denmark, Czechia, "
       "Bulgaria… and more.", 12.5, GREY, False)]],
    line_spacing=1.12)
txt(s, Inches(7.35), Inches(5.6), Inches(5.2), Inches(1.0),
    [[("Note: ", 12.5, GOLD, True),
      ("working papers concentrate the volume in ~13 large banks; the smaller "
       "ones have (for now) only speeches.", 12.5, GREY, False)]],
    line_spacing=1.1)
footer(s, 4)

# ===========================================================================
# 5. THE 5 MISSING COUNTRIES
# ===========================================================================
s = slide(); header(s, "Coverage", "The 2 countries still missing")
txt(s, Inches(0.6), Inches(1.65), Inches(12), Inches(0.7),
    [[("61 of the 63 target banks are present. The 3 previously 'missing' "
       "(Latvia, Lithuania, Luxembourg) were recovered — their absence was an "
       "attribution bug, now fixed. Two genuine absences remain:", 16, GREY, False)]])
miss = [("Peru — Central Reserve Bank", "Publishes mainly in Spanish — no EN version indexed by the BIS"),
        ("Vietnam — State Bank", "Publishes mainly in Vietnamese — same")]
y = Inches(2.9); rh = Inches(0.9)
for i, (lab, why) in enumerate(miss):
    rect(s, Inches(0.6), y, Inches(12.1), rh-Inches(0.12), LIGHT)
    rect(s, Inches(0.6), y, Inches(0.1), rh-Inches(0.12), GOLD)
    txt(s, Inches(0.9), y+Inches(0.06), Inches(4.6), Inches(0.6),
        [[(lab, 15, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.6), y+Inches(0.06), Inches(6.9), Inches(0.6),
        [[(why, 13.5, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += rh
footer(s, 5)

# ===========================================================================
# 6. ROADMAP
# ===========================================================================
s = slide(); header(s, "What's next", "What remains to collect")
txt(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.5),
    [[("✅ Speeches and research papers are collected. Three workstreams remain, "
       "from highest priority to most complex:", 15, GREY, False)]])
items = [
    ("1", "Decisions — major banks done", GOLD,
     "Rate decisions, policy statements & meeting minutes: 1,872 collected "
     "(ECB, Fed, Australia, Bank of England).",
     "Largely done for the majors; extending site-by-site to the rest."),
    ("2", "Press conferences", BLUE,
     "The transcripts: the 'why' behind each decision, journalists' questions "
     "included.",
     "Most demanding technically. Starting with the Fed and the ECB."),
    ("3", "Financial-stability & annual reports", NAVY,
     "Banks' financial-stability reports and annual reports.",
     "Simple patterns to add per bank."),
]
y = Inches(2.35); rh = Inches(1.35)
for num, title, col, body, status in items:
    rect(s, Inches(0.6), y, Inches(12.1), rh-Inches(0.13), LIGHT)
    rect(s, Inches(0.6), y, Inches(1.0), rh-Inches(0.13), col)
    txt(s, Inches(0.6), y, Inches(1.0), rh-Inches(0.13),
        [[(num, 34, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.8), y+Inches(0.12), Inches(10.7), Inches(0.4),
        [[(title, 17, NAVY, True)]])
    txt(s, Inches(1.8), y+Inches(0.52), Inches(10.7), Inches(0.5),
        [[(body, 13, GREY, False)]], line_spacing=1.0)
    txt(s, Inches(1.8), y+rh-Inches(0.42), Inches(10.7), Inches(0.35),
        [[("Status: ", 12, col, True), (status, 12, GREY, False)]])
    y += rh
footer(s, 6)

# ===========================================================================
# 7. SUMMARY
# ===========================================================================
s = slide(); header(s, "Summary", "Today → goal")
cols = [("", "Today", "Goal"),
        ("Documents", "38,413", "45,000+"),
        ("Countries covered", "61 / 63", "63 / 63"),
        ("Document types", "speeches + research + decisions + minutes + reports", "+ press conferences"),
        ("Source nature", "100% official", "unchanged — that's the principle")]
y = Inches(1.8); rh = Inches(0.72)
cw = [Inches(3.6), Inches(4.25), Inches(4.25)]
cx0 = Inches(0.6)
for i, row in enumerate(cols):
    header_row = (i == 0)
    cx = cx0
    for j, cell in enumerate(row):
        if header_row:
            rect(s, cx, y, cw[j], rh, NAVY if j == 0 else (BLUE if j == 1 else GOLD))
            color = WHITE
        else:
            rect(s, cx, y, cw[j], rh, LIGHT if i % 2 else WHITE)
            color = NAVY if j == 0 else GREY
        bold = header_row or j == 0
        txt(s, cx+Inches(0.2), y, cw[j]-Inches(0.3), rh,
            [[(cell, 14.5 if not header_row else 15, color, bold)]],
            anchor=MSO_ANCHOR.MIDDLE)
        cx += cw[j]
    y += rh

rect(s, Inches(0.6), y+Inches(0.2), Inches(12.1), Inches(1.15), NAVY)
txt(s, Inches(0.9), y+Inches(0.35), Inches(11.5), Inches(0.9),
    [[("What will never change: ", 15, GOLD, True),
      ("first-hand official sources only, traceable, with no commentary or "
       "generated content. A reliable, exhaustive knowledge base of the "
       "world's central banks' own words.",
       15, WHITE, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.08)
footer(s, 7)

prs.save("PRESENTATION_EN.pptx")
print("OK — PRESENTATION_EN.pptx,", len(prs.slides._sldIdLst), "slides")
